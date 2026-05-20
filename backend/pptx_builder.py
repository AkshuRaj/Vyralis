import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')

# Color scheme
NAVY = RGBColor(10, 35, 66)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(230, 57, 70)
GRAY = RGBColor(107, 114, 128)
LIGHT_GRAY = RGBColor(243, 244, 246)
LIGHT_BLUE = RGBColor(220, 240, 255)
ORANGE = RGBColor(255, 153, 0)
GREEN = RGBColor(34, 197, 94)

# Slide dimensions
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Margins
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
CONTENT_TOP = Inches(1.1)


def normalize_analysis(analysis: dict) -> dict:
    """Normalize Groq response keys to match what pptx_builder expects."""

    if "overall_winner" in analysis and "winner" not in analysis:
        analysis["winner"] = analysis["overall_winner"]

    if "strategic_recommendations" in analysis and "recommendations" not in analysis:
        analysis["recommendations"] = analysis["strategic_recommendations"]

    if "performance_scorecard" in analysis and "scorecard" not in analysis:
        analysis["scorecard"] = analysis["performance_scorecard"]

    if "competitive_gap_analysis" in analysis and "gap_analysis" not in analysis:
        gap = analysis["competitive_gap_analysis"]
        analysis["gap_analysis"] = {
            "untapped_topics": [
                gap.get("content_gap", ""),
                gap.get("biggest_opportunity", ""),
                gap.get("audience_expectation", "")
            ],
            "opportunity": gap.get("supporting_evidence", gap.get("biggest_opportunity", "")),
            "format_gaps": []
        }

    if "posting_strategy" in analysis and "posting_analysis" not in analysis:
        ps = analysis["posting_strategy"]
        analysis["posting_analysis"] = {
            "most_consistent": ps.get("most_consistent", ""),
            "insights": ps.get("frequency_insight", "") + " " + ps.get("posting_pattern", "")
        }

    ea = analysis.get("engagement_analysis", {})
    if "insights" not in ea:
        ea["insights"] = ea.get("engagement_rate_insight", "")
    if "engagement_tips" not in ea:
        ea["engagement_tips"] = ea.get("engagement_opportunities", [])

    scorecard = analysis.get("scorecard", {})
    for company, scores in scorecard.items():
        if not isinstance(scores, dict):
            continue
        if "subscriber_score" not in scores:
            scores["subscriber_score"] = scores.get("subscriber_growth_potential", 5)
        if "engagement_score" not in scores:
            scores["engagement_score"] = scores.get("engagement_quality", 5)
        if "consistency_score" not in scores:
            scores["consistency_score"] = scores.get("content_consistency", 5)
        if "content_quality_score" not in scores:
            scores["content_quality_score"] = scores.get("strategic_positioning", 5)
        if "overall_score" not in scores:
            scores["overall_score"] = scores.get("overall_score", 5)

    return analysis


def add_slide_header(slide, title_text):
    """Add consistent navy header bar to every slide."""
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.85)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.color.rgb = NAVY

    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(0.15),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_accent_bar(slide):
    """Add thin accent line below header."""
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0.85),
        SLIDE_WIDTH, Inches(0.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.color.rgb = ACCENT


def build_pptx(company_data: list, analysis: dict, failed_companies: list = None) -> bytes:
    """Build professional PowerPoint report."""

    analysis = normalize_analysis(analysis)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    company_names = [c.get("company_name") for c in company_data if "error" not in c]

    slide_cover(prs, company_names, analysis)
    slide_executive_summary(prs, analysis, len(company_data))
    slide_channel_overview(prs, company_data)
    slide_subscriber_chart(prs, company_data)
    slide_top_videos(prs, company_data)
    slide_content_themes(prs, analysis)
    slide_posting_frequency(prs, company_data, analysis)
    slide_engagement_chart(prs, company_data, analysis)
    slide_gap_analysis(prs, analysis)
    slide_recommendations(prs, analysis)
    slide_scorecard(prs, analysis)
    slide_thank_you(prs, analysis)

    if failed_companies and len(failed_companies) > 0:
        slide_notes(prs, failed_companies)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def slide_cover(prs, company_names, analysis):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Accent bar at bottom
    bottom_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(5.2),
        SLIDE_WIDTH, Inches(0.15)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT
    bottom_bar.line.color.rgb = ACCENT

    # Report label
    label_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(0.8),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.4)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "VIDEO COMPETITOR INTELLIGENCE REPORT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Company names
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.4),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = " vs ".join(company_names)
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(3.1),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "YouTube Channel & Video Marketing Analysis"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Winner badge
    winner = analysis.get("winner", "")
    if winner:
        badge_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(3.8),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        tf = badge_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🏆  Current Leader: {winner}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(4.9),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.4)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def slide_executive_summary(prs, analysis, total_companies):
    """Slide 2: Executive Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Executive Summary")
    add_accent_bar(slide)

    # Summary text box with light background
    summary_bg = slide.shapes.add_shape(
        1, MARGIN_LEFT, CONTENT_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1.8)
    )
    summary_bg.fill.solid()
    summary_bg.fill.fore_color.rgb = LIGHT_GRAY
    summary_bg.line.color.rgb = LIGHT_GRAY

    summary_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.15),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.4), Inches(1.5)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = analysis.get("executive_summary", "No summary available.")
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.5

    # Stats row
    companies_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(3.1),
        Inches(3), Inches(0.8)
    )
    tf = companies_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(total_companies)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Companies Analysed"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Winner box
    winner = analysis.get("winner", "")
    reasoning = analysis.get("channel_comparison", {}).get("reasoning", "")

    winner_box = slide.shapes.add_textbox(
        Inches(3.5), Inches(3.1),
        Inches(6), Inches(1.8)
    )
    tf = winner_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"🏆  Leader: {winner}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = reasoning[:180] if reasoning else ""
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3


def slide_channel_overview(prs, company_data):
    """Slide 3: Channel Overview Table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Channel Overview Comparison")
    add_accent_bar(slide)

    rows = len(company_data) + 1
    cols = 5
    table_shape = slide.shapes.add_table(
        rows, cols,
        MARGIN_LEFT, CONTENT_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(3.8)
    )
    table = table_shape.table

    headers = ["Channel", "Subscribers", "Total Videos", "Total Views", "Country"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    max_subs = max(
        [c.get("channel_stats", {}).get("subscriber_count", 0) for c in company_data],
        default=0
    )

    for row_idx, company in enumerate(company_data, 1):
        stats = company.get("channel_stats", {})
        is_max = stats.get("subscriber_count", 0) == max_subs
        bg_color = LIGHT_BLUE if is_max else (LIGHT_GRAY if row_idx % 2 == 0 else WHITE)

        data = [
            company.get("company_name", ""),
            f"{stats.get('subscriber_count', 0):,}",
            f"{stats.get('video_count', 0):,}",
            f"{stats.get('view_count', 0):,}",
            stats.get("country", "N/A"),
        ]

        for col_idx, value in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            tf = cell.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.bold = is_max
            tf.paragraphs[0].font.color.rgb = NAVY if is_max else GRAY
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def slide_subscriber_chart(prs, company_data):
    """Slide 4: Subscriber Comparison Chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Subscriber Count Comparison")
    add_accent_bar(slide)

    companies = [c.get("company_name", "") for c in company_data]
    subscribers = [c.get("channel_stats", {}).get("subscriber_count", 0) for c in company_data]

    fig, ax = plt.subplots(figsize=(9, 3.8), facecolor="white")
    colors = ["#E63946" if i == 0 else "#0A2342" for i in range(len(companies))]
    bars = ax.barh(companies, subscribers, color=colors, height=0.5)

    ax.set_xlabel("Subscriber Count", fontsize=10, color="#6B7280")
    ax.tick_params(colors="#6B7280", labelsize=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E5E7EB")
    ax.spines["bottom"].set_color("#E5E7EB")
    ax.set_facecolor("white")

    for bar, val in zip(bars, subscribers):
        ax.text(
            val + max(subscribers) * 0.01,
            bar.get_y() + bar.get_height() / 2,
            f"{val:,}",
            va="center", fontsize=9, color="#6B7280"
        )

    fig.tight_layout()

    img_bytes = io.BytesIO()
    fig.savefig(img_bytes, format="png", dpi=150, bbox_inches="tight")
    img_bytes.seek(0)
    plt.close(fig)

    slide.shapes.add_picture(
        img_bytes, MARGIN_LEFT, CONTENT_TOP,
        width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    )


def slide_top_videos(prs, company_data):
    """Slide 5: Top Performing Videos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Top Performing Video per Company")
    add_accent_bar(slide)

    current_top = CONTENT_TOP
    row_height = Inches(1.05)

    for company in company_data:
        videos = company.get("top_videos", [])
        if not videos:
            continue

        top_video = videos[0]

        # Card background
        card = slide.shapes.add_shape(
            1, MARGIN_LEFT, current_top,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, row_height - Inches(0.05)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(229, 231, 235)

        # Company label
        label = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.15), current_top + Inches(0.08),
            Inches(2), Inches(0.3)
        )
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = company.get("company_name", "")
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        # Video title
        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.15), current_top + Inches(0.4),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.3), Inches(0.35)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        title = top_video.get("title", "")[:70]
        p.text = f"🎬  {title}"
        p.font.size = Pt(10)
        p.font.color.rgb = NAVY

        # Stats
        stats_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.15), current_top + Inches(0.72),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.25)
        )
        tf = stats_box.text_frame
        p = tf.paragraphs[0]
        views = top_video.get("view_count", 0)
        likes = top_video.get("like_count", 0)
        p.text = f"👁  {views:,} views   👍  {likes:,} likes"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY

        current_top += row_height


def slide_content_themes(prs, analysis):
    """Slide 6: Content Themes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Content Themes & Topics")
    add_accent_bar(slide)

    content_themes = analysis.get("content_themes", {})
    companies_list = list(content_themes.keys())[:4]

    col_count = 2 if len(companies_list) > 2 else len(companies_list)
    col_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.3)) / col_count

    for idx, company in enumerate(companies_list):
        themes = content_themes.get(company, {})
        col = idx % 2
        row = idx // 2
        left = MARGIN_LEFT + col * (col_width + Inches(0.3))
        top = CONTENT_TOP + row * Inches(2.1)

        # Company header
        header_bg = slide.shapes.add_shape(
            1, left, top, col_width, Inches(0.32)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = NAVY
        header_bg.line.color.rgb = NAVY

        header_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.04),
            col_width - Inches(0.2), Inches(0.28)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = company
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Topics
        topics_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.38),
            col_width - Inches(0.2), Inches(1.6)
        )
        tf = topics_box.text_frame
        tf.word_wrap = True

        topics = themes.get("main_topics", [])
        first = True
        for topic in topics[:4]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"• {topic}"
            p.font.size = Pt(10)
            p.font.color.rgb = GRAY

        # Missing topics
        missing = themes.get("missing_topics", [])
        if missing:
            p = tf.add_paragraph()
            p.text = f"Gap: {missing[0]}"
            p.font.size = Pt(9)
            p.font.color.rgb = ACCENT
            p.font.italic = True


def slide_posting_frequency(prs, company_data, analysis):
    """Slide 7: Posting Frequency"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Posting Frequency & Consistency")
    add_accent_bar(slide)

    most_consistent = analysis.get("posting_analysis", {}).get("most_consistent", "")

    rows = len(company_data) + 1
    table_shape = slide.shapes.add_table(
        rows, 3,
        MARGIN_LEFT, CONTENT_TOP,
        Inches(5), Inches(0.45 * rows)
    )
    table = table_shape.table

    for col_idx, header in enumerate(["Company", "Posts/Month", "Status"]):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

    for row_idx, company in enumerate(company_data, 1):
        name = company.get("company_name", "")
        freq = company.get("posting_frequency", 0)
        is_best = name == most_consistent
        bg_color = LIGHT_BLUE if is_best else (LIGHT_GRAY if row_idx % 2 == 0 else WHITE)

        for col_idx, value in enumerate([
            name,
            str(freq),
            "🏆 Most Consistent" if is_best else "Regular"
        ]):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            tf = cell.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.bold = is_best
            tf.paragraphs[0].font.color.rgb = NAVY if is_best else GRAY

    # Insights
    insights = analysis.get("posting_analysis", {}).get("insights", "")
    if insights:
        insights_box = slide.shapes.add_textbox(
            MARGIN_LEFT, CONTENT_TOP + Inches(0.45 * rows) + Inches(0.2),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(2)
        )
        tf = insights_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insights
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.line_spacing = 1.4


def slide_engagement_chart(prs, company_data, analysis):
    """Slide 8: Engagement Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Engagement Analysis")
    add_accent_bar(slide)

    companies = [c.get("company_name", "") for c in company_data]
    avg_views = [c.get("averages", {}).get("avg_views", 0) for c in company_data]
    avg_likes = [c.get("averages", {}).get("avg_likes", 0) for c in company_data]
    avg_comments = [c.get("averages", {}).get("avg_comments", 0) for c in company_data]

    fig, ax = plt.subplots(figsize=(8, 3), facecolor="white")
    x = range(len(companies))
    width = 0.25

    ax.bar([i - width for i in x], avg_views, width, label="Avg Views", color="#E63946")
    ax.bar(list(x), avg_likes, width, label="Avg Likes", color="#0A2342")
    ax.bar([i + width for i in x], avg_comments, width, label="Avg Comments", color="#6B7280")

    ax.set_xticks(list(x))
    ax.set_xticklabels(companies, fontsize=10)
    ax.set_ylabel("Count", fontsize=9, color="#6B7280")
    ax.legend(fontsize=9)
    ax.tick_params(colors="#6B7280")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_facecolor("white")

    fig.tight_layout()

    img_bytes = io.BytesIO()
    fig.savefig(img_bytes, format="png", dpi=150, bbox_inches="tight")
    img_bytes.seek(0)
    plt.close(fig)

    slide.shapes.add_picture(
        img_bytes, MARGIN_LEFT, CONTENT_TOP,
        width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.5)
    )

    insights = analysis.get("engagement_analysis", {}).get("insights", "")
    if insights:
        ins_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(4.2),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
        )
        tf = ins_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = insights
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY


def slide_gap_analysis(prs, analysis):
    """Slide 9: Gap Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Opportunities You Are Missing")
    add_accent_bar(slide)

    gap = analysis.get("gap_analysis", {})

    # Left column — untapped topics
    left_box = slide.shapes.add_textbox(
        MARGIN_LEFT, CONTENT_TOP, Inches(4.5), Inches(4)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Untapped Content Topics"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY

    for topic in gap.get("untapped_topics", [])[:5]:
        if topic:
            p = tf.add_paragraph()
            p.text = f"▶  {topic}"
            p.font.size = Pt(11)
            p.font.color.rgb = GRAY
            p.line_spacing = 1.4

    # Right column — opportunity
    right_bg = slide.shapes.add_shape(
        1, Inches(5.2), CONTENT_TOP,
        Inches(4.3), Inches(4)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(10, 35, 66)
    right_bg.line.color.rgb = ACCENT

    right_box = slide.shapes.add_textbox(
        Inches(5.4), CONTENT_TOP + Inches(0.15),
        Inches(3.9), Inches(3.7)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Biggest Opportunity"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = gap.get("opportunity", "")
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.line_spacing = 1.4


def slide_recommendations(prs, analysis):
    """Slide 10: Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Strategic Recommendations")
    add_accent_bar(slide)

    recommendations = analysis.get("recommendations", [])[:5]
    current_top = CONTENT_TOP

    for rec in recommendations:
        priority = rec.get("priority", "Medium")

        if priority in ("Immediate", "High"):
            priority_color = ACCENT
        elif priority == "Medium":
            priority_color = ORANGE
        else:
            priority_color = GREEN

        card_height = Inches(0.88)

        # Card background
        card = slide.shapes.add_shape(
            1, MARGIN_LEFT, current_top,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = priority_color
        card.line.width = Pt(2)

        # Priority badge
        badge = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.15), current_top + Inches(0.08),
            Inches(1.2), Inches(0.28)
        )
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = f"● {priority}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = priority_color

        # Action
        action_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(1.4), current_top + Inches(0.06),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(1.6), Inches(0.38)
        )
        tf = action_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rec.get("action", "")
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Impact
        impact_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(1.4), current_top + Inches(0.48),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(1.6), Inches(0.32)
        )
        tf = impact_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Expected: {rec.get('expected_impact', '')}"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY
        p.font.italic = True

        current_top += card_height + Inches(0.05)


def slide_scorecard(prs, analysis):
    """Slide 11: Scorecard"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Performance Scorecard")
    add_accent_bar(slide)

    scorecard = analysis.get("scorecard", {})
    companies = list(scorecard.keys())
    winner = analysis.get("winner", analysis.get("overall_winner", ""))

    if not companies:
        return

    rows = len(companies) + 1
    cols = 6
    table_shape = slide.shapes.add_table(
        rows, cols,
        MARGIN_LEFT, CONTENT_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(3.2)
    )
    table = table_shape.table

    headers = ["Company", "Subscribers", "Engagement", "Consistency", "Content", "Overall /10"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.text = header
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    for row_idx, company in enumerate(companies, 1):
        scores = scorecard.get(company, {})
        if not isinstance(scores, dict):
            continue

        is_winner = company == winner
        bg_color = ACCENT if is_winner else (LIGHT_GRAY if row_idx % 2 == 0 else WHITE)
        text_color = WHITE if is_winner else NAVY

        data = [
            company,
            str(scores.get("subscriber_score", "-")),
            str(scores.get("engagement_score", "-")),
            str(scores.get("consistency_score", "-")),
            str(scores.get("content_quality_score", "-")),
            str(scores.get("overall_score", "-")),
        ]

        for col_idx, value in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            tf = cell.text_frame
            tf.text = value
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.bold = is_winner
            tf.paragraphs[0].font.color.rgb = text_color
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Key insight below table
    key_insight = analysis.get("key_insight", "")
    if key_insight:
        insight_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(4.5),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.8)
        )
        tf = insight_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"💡  {key_insight}"
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
        p.font.italic = True


def slide_thank_you(prs, analysis):
    """Slide 12: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    # Accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(2.5),
        SLIDE_WIDTH, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.color.rgb = ACCENT

    winner = analysis.get("winner", "")
    winner_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(0.8),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1.5)
    )
    tf = winner_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"🏆  {winner} leads the competition"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    summary = analysis.get("executive_summary", "")
    if summary:
        first_sentence = summary.split(".")[0] + "."
        sum_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(2.8),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1.2)
        )
        tf = sum_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = first_sentence
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    footer_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(4.8),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.4)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Generated by Video Competitor Intelligence Tool"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER


def slide_notes(prs, failed_companies):
    """Optional slide: Data notes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    add_slide_header(slide, "Data Collection Notes")
    add_accent_bar(slide)

    info_box = slide.shapes.add_textbox(
        MARGIN_LEFT, CONTENT_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(3.5)
    )
    tf = info_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Companies Unable to Locate:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    for failed in failed_companies:
        company_name = failed.get("company", "Unknown") if isinstance(failed, dict) else str(failed)
        reason = failed.get("reason", "No YouTube channel found") if isinstance(failed, dict) else "No YouTube channel found"
        p = tf.add_paragraph()
        p.text = f"• {company_name}: {reason}"
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY

    p = tf.add_paragraph()
    p.text = "\nThe analysis is based on companies for which official YouTube channels were found."
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)